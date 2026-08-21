#!/usr/bin/env python3
"""Generate B2B sales deck for Королевская Биржа (event managers)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Brand palette
BG = RGBColor(0x12, 0x1A, 0x2E)
BG_ALT = RGBColor(0x1A, 0x27, 0x44)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
WHITE = RGBColor(0xF5, 0xF5, 0xF0)
MUTED = RGBColor(0xA8, 0xB2, 0xC4)
ACCENT = RGBColor(0x4A, 0x90, 0xD9)
RED_MUTED = RGBColor(0xE8, 0x8A, 0x8A)
GREEN_MUTED = RGBColor(0x8A, 0xD4, 0xA8)

OUTPUT = Path(__file__).resolve().parent.parent / "Королевская_Биржа_для_ивент-агентств.pptx"


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        pass  # python-pptx limited alpha on shapes
    return shape


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size=14,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.LEFT,
    font_name="Calibri",
    line_spacing=1.15,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return box


def add_bullets(slide, left, top, width, height, items, *, font_size=13, color=WHITE, title=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = title
        p.space_after = Pt(6)
        r = p.runs[0]
        r.font.size = Pt(font_size + 1)
        r.font.bold = True
        r.font.color.rgb = GOLD
        r.font.name = "Calibri"
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.space_after = Pt(4)
        p.line_spacing = 1.1
        r = p.runs[0]
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return box


def add_footer(slide, text: str):
    add_text_box(
        slide,
        Inches(0.5),
        Inches(7.05),
        Inches(12.3),
        Inches(0.35),
        text,
        font_size=9,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def slide_header(slide, title: str, subtitle: str = ""):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.12), GOLD)
    add_text_box(
        slide,
        Inches(0.55),
        Inches(0.35),
        Inches(12.2),
        Inches(0.9),
        title,
        font_size=28,
        bold=True,
        color=WHITE,
    )
    if subtitle:
        add_text_box(
            slide,
            Inches(0.55),
            Inches(1.15),
            Inches(12.2),
            Inches(0.55),
            subtitle,
            font_size=15,
            color=MUTED,
        )


def build_slide_1(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    slide_header(
        slide,
        "Вовлечь всю команду за один блок программы",
        "Королевская Биржа — экономическая командная игра для корпоративов и бизнес-ивентов",
    )

    # Left column label
    add_text_box(
        slide, Inches(0.55), Inches(1.85), Inches(5.8), Inches(0.35),
        "Запрос рынка", font_size=12, bold=True, color=RED_MUTED,
    )
    add_bullets(
        slide, Inches(0.55), Inches(2.2), Inches(5.85), Inches(2.8),
        [
            "Участники остаются зрителями, а не в игре",
            "Классический тимбилдинг быстро надоедает",
            "Сложно удержать внимание 30–200 человек в зале",
            "Нужен формат, который выглядит серьёзно для руководства",
        ],
        font_size=13,
        color=MUTED,
    )

    # Right column
    add_text_box(
        slide, Inches(6.9), Inches(1.85), Inches(5.8), Inches(0.35),
        "Наше решение", font_size=12, bold=True, color=GREEN_MUTED,
    )
    add_bullets(
        slide, Inches(6.9), Inches(2.2), Inches(5.85), Inches(3.2),
        [
            "«Королевская Биржа» — стратегия с торгами, инвестициями и рейтингом в реальном времени",
            "Каждый играет с телефона; общий экран — для драматургии и соревнования",
            "5–30 игроков · 10 раундов · 60–90 минут",
            "Сюжет, события рынка, видео-введение — погружение без реквизита",
        ],
        font_size=13,
        color=WHITE,
    )

    # Bottom promise bar
    add_rect(slide, Inches(0.55), Inches(5.85), Inches(12.25), Inches(0.75), BG_ALT)
    add_rect(slide, Inches(0.55), Inches(5.85), Inches(0.08), Inches(0.75), GOLD)
    add_text_box(
        slide, Inches(0.85), Inches(6.0), Inches(11.7), Inches(0.5),
        "Готовый digital-формат: вы продаёте клиенту впечатление — мы даём технологию и сценарий.",
        font_size=14,
        bold=True,
        color=GOLD,
        align=PP_ALIGN.LEFT,
    )
    add_footer(slide, "Королевская Биржа · B2B-формат для ивент-агентств")


def build_slide_2(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    slide_header(slide, "Как это работает на площадке", "Turnkey-формат для event-команды")

    steps = [
        ("1", "Старт", "Ведущий открывает игру на проекторе.\nУчастники заходят в Telegram Mini App по коду."),
        ("2", "Раунды", "Покупка и продажа ресурсов, строительство объектов,\nреакция на события рынка."),
        ("3", "Финал", "Турнирная таблица, победитель, награждение —\nсильный эмоциональный пик программы."),
    ]
    x0 = Inches(0.55)
    w = Inches(3.95)
    for i, (num, title, body) in enumerate(steps):
        left = x0 + i * (w + Inches(0.2))
        add_rect(slide, left, Inches(1.95), w, Inches(1.85), BG_ALT)
        add_text_box(slide, left + Inches(0.15), Inches(2.05), Inches(0.5), Inches(0.45),
                     num, font_size=22, bold=True, color=GOLD)
        add_text_box(slide, left + Inches(0.55), Inches(2.05), Inches(3.2), Inches(0.4),
                     title, font_size=16, bold=True, color=WHITE)
        add_text_box(slide, left + Inches(0.15), Inches(2.5), Inches(3.65), Inches(1.2),
                     body, font_size=11, color=MUTED, line_spacing=1.2)

    add_text_box(
        slide, Inches(0.55), Inches(4.05), Inches(12), Inches(0.35),
        "Почему event-агентствам это удобно",
        font_size=14, bold=True, color=GOLD,
    )

    benefits_left = [
        "Не нужны реквизит и раздаточные материалы",
        "Проектор + звук; участники не в очереди к одному столу",
        "Админ-панель: раунды, откат, настройка событий",
    ]
    benefits_right = [
        "Синхронизация в реальном времени (WebSocket)",
        "Стратегия и риск — понятно топ-менеджменту",
        "Гибкость: брендинг, длительность, корпоративная легенда",
    ]
    add_bullets(slide, Inches(0.55), Inches(4.45), Inches(6.0), Inches(2.0), benefits_left, font_size=12)
    add_bullets(slide, Inches(6.75), Inches(4.45), Inches(6.0), Inches(2.0), benefits_right, font_size=12)

    add_rect(slide, Inches(0.55), Inches(6.55), Inches(12.25), Inches(0.42), BG_ALT)
    add_text_box(
        slide, Inches(0.7), Inches(6.58), Inches(12.0), Inches(0.38),
        "5–30 игроков  ·  10 раундов  ·  9 ресурсов  ·  11 объектов  ·  "
        "Проектор + Telegram Mini App  ·  Видео-введение  ·  Облачный деплой",
        font_size=10,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "Королевская Биржа · Технология под ключ")


def build_slide_3(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    slide_header(
        slide,
        "Добавьте в линейку формат, который клиенты запомнят",
        "Партнёрское предложение для ивент-менеджеров и event-компаний",
    )

    tags = "Корпоративы и тимбилдинги  ·  Конференции и business days  ·  Камерные встречи руководителей"
    add_text_box(slide, Inches(0.55), Inches(1.85), Inches(12.2), Inches(0.4), tags, font_size=12, color=MUTED)

    # Two partnership cards
    cards = [
        (
            "White-label для агентства",
            [
                "Вы продаёте клиенту — мы обеспечиваем платформу",
                "Методичка и сценарий для ведущего",
                "Демо и обучение вашей команде",
            ],
        ),
        (
            "Под ключ",
            [
                "Бриф → техпроверка → игровой день",
                "Ваш ведущий или наш — на выбор",
                "Кастомизация сюжета под отрасль клиента",
            ],
        ),
    ]
    for i, (title, items) in enumerate(cards):
        left = Inches(0.55) + i * Inches(6.35)
        add_rect(slide, left, Inches(2.35), Inches(6.1), Inches(2.15), BG_ALT)
        add_bullets(slide, left + Inches(0.2), Inches(2.45), Inches(5.7), Inches(2.0), items, title=title, font_size=12)

    add_text_box(
        slide, Inches(0.55), Inches(4.65), Inches(12), Inches(0.35),
        "Что получает ивент-менеджер",
        font_size=13, bold=True, color=GOLD,
    )
    add_bullets(
        slide, Inches(0.55), Inches(5.0), Inches(12.2), Inches(1.0),
        [
            "Готовый продукт в каталоге с понятным таймингом и КП",
            "Меньше риска провала активности в зале · сильный визуал для отчёта заказчику",
        ],
        font_size=12,
        color=MUTED,
    )

    # CTA block
    add_rect(slide, Inches(0.55), Inches(5.95), Inches(12.25), Inches(1.05), GOLD)
    add_text_box(
        slide, Inches(0.75), Inches(6.05), Inches(11.9), Inches(0.45),
        "Запишитесь на демо 20 минут",
        font_size=20, bold=True, color=BG, align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide, Inches(0.75), Inches(6.48), Inches(11.9), Inches(0.4),
        "[ Имя ]  ·  [ телефон / Telegram ]  ·  [ email ]  ·  [ сайт ]",
        font_size=12, color=BG, align=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide, Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.3),
        "Техподдержка в день мероприятия  ·  Тестовый прогон до ивента  ·  Условия — по запросу",
        font_size=9, color=MUTED, align=PP_ALIGN.CENTER,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
